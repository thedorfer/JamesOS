from pathlib import Path
from datetime import datetime
import tempfile
import zipfile

from jamesos.config import VAULT

ARCHIVE_ATTACHMENTS = VAULT / "Archive" / "Attachments"
KNOWLEDGE_FILES = VAULT / "JamesOS" / "Knowledge" / "Files"


def _extract_pdf(path: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    return "\n\n".join(
        f"## Page {i}\n{page.extract_text() or ''}"
        for i, page in enumerate(reader.pages, start=1)
    )


def _extract_docx(path: Path) -> str:
    import docx
    doc = docx.Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_html(path: Path) -> str:
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="ignore"), "html.parser")
    return soup.get_text("\n")


def _extract_rtf(path: Path) -> str:
    from striprtf.striprtf import rtf_to_text
    return rtf_to_text(path.read_text(encoding="utf-8", errors="ignore"))


def _extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(str(path), read_only=True, data_only=True)
    lines = []
    for ws in wb.worksheets:
        lines.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            values = [str(v) for v in row if v is not None]
            if values:
                lines.append(" | ".join(values))
    return "\n".join(lines)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    lines = []
    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"# Slide {i}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
    return "\n".join(lines)


def _extract_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def extract_file_text(path: Path) -> str:
    ext = path.suffix.lower()

    if ext == ".pdf":
        return _extract_pdf(path)
    if ext == ".docx":
        return _extract_docx(path)
    if ext in {".txt", ".md", ".csv", ".json", ".xml", ".sql", ".py", ".log", ".yaml", ".yml"}:
        return _extract_text(path)
    if ext in {".html", ".htm"}:
        return _extract_html(path)
    if ext == ".rtf":
        return _extract_rtf(path)
    if ext == ".xlsx":
        return _extract_xlsx(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if ext == ".zip":
        return _extract_zip(path)

    return ""


def _extract_zip(path: Path) -> str:
    lines = [f"# ZIP Archive: {path.name}", ""]
    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        with zipfile.ZipFile(path) as z:
            z.extractall(tmp_path)

        for child in sorted(tmp_path.rglob("*")):
            if not child.is_file():
                continue
            text = extract_file_text(child)
            if text.strip():
                lines.append(f"## File: {child.relative_to(tmp_path)}")
                lines.append(text[:15000])
                lines.append("")
    return "\n".join(lines)


def build_file_knowledge() -> str:
    KNOWLEDGE_FILES.mkdir(parents=True, exist_ok=True)

    processed = 0
    skipped = 0

    for path in ARCHIVE_ATTACHMENTS.rglob("*"):
        if not path.is_file():
            continue

        text = extract_file_text(path)
        if not text.strip():
            skipped += 1
            continue

        rel = path.relative_to(VAULT).as_posix()
        safe_name = "".join(c if c.isalnum() or c in " ._-" else "_" for c in path.stem)
        note = KNOWLEDGE_FILES / f"{safe_name}.md"

        note.write_text(f"""# {path.stem}

Type: file_knowledge
Source File: [[{Path(rel).with_suffix('').as_posix()}]]
Original Path: {rel}
Extracted At: {datetime.now().strftime('%Y-%m-%d %H:%M')}

## Extracted Text

{text[:50000]}
""", encoding="utf-8")

        processed += 1

    return f"File intelligence complete. Processed: {processed}. Skipped: {skipped}."
